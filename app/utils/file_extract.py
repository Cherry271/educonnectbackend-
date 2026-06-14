def extract_text_from_file(content: bytes, filename: str) -> str:
    ext = filename.lower().split(".")[-1] if "." in filename else ""
    try:
        if ext == "pdf":
            import io

            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        if ext == "docx":
            import io

            from docx import Document

            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        if ext == "pptx":
            import io

            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        if ext in ("txt", "md"):
            return content.decode("utf-8", errors="ignore")
    except Exception:
        pass
    return ""
